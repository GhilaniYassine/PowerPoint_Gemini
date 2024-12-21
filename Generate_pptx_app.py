import google.generativeai as genai
from pptx import Presentation
from pptx.util import Inches, Pt
import os
import streamlit as st
import base64
from dotenv import load_dotenv

load_dotenv()

# Define custom formatting options
TITLE_FONT_SIZE = Pt(30)
SLIDE_FONT_SIZE = Pt(16)

def generate_slide_titles(topic):
    """
    Generate slide titles for a given topic using Gemini API.
    Assumes the API key is loaded and accessible from environment variables.
    """
    # Configure the Gemini API
    genai.configure(api_key=os.getenv("GEMINI_API_KEY"))

    # Initialize the model
    model = genai.GenerativeModel("gemini-1.5-flash")

    # Generate slide titles
    prompt = f"Generate 5 slide titles for the topic '{topic}'."
    response = model.generate_content(prompt)

    # Process and return the response
    return response.text.split("\n")

def generate_slide_content(slide_title):
    """
    Generate detailed content for a given slide title using Gemini API.
    Assumes the API key is loaded and accessible from environment variables.
    """
    # Configure the Gemini API
    genai.configure(api_key=os.getenv("GEMINI_API_KEY"))

    # Initialize the model
    model = genai.GenerativeModel("gemini-1.5-flash")

    # Generate slide content
    prompt = f"Generate content of 6 lines of information based on the slide title: '{slide_title}'."
    response = model.generate_content(prompt)

    # Process and return the response
    return response.text.strip()

def create_presentation(topic, slide_titles, slide_contents):
    prs = Presentation()
    slide_layout = prs.slide_layouts[1]

    title_slide = prs.slides.add_slide(prs.slide_layouts[0])
    title_slide.shapes.title.text = topic

    for slide_title, slide_content in zip(slide_titles, slide_contents):
        slide = prs.slides.add_slide(slide_layout)
        slide.shapes.title.text = slide_title
        slide.shapes.placeholders[1].text = slide_content

        # Customize font size for titles and content
        slide.shapes.title.text_frame.paragraphs[0].font.size = TITLE_FONT_SIZE
        for shape in slide.shapes:
            if shape.has_text_frame:
                text_frame = shape.text_frame
                for paragraph in text_frame.paragraphs:
                    paragraph.font.size = SLIDE_FONT_SIZE

    os.makedirs("generated_ppt", exist_ok=True)
    prs.save(f"generated_ppt/{topic}_presentation.pptx")

def main():
    # Add a custom logo
    st.image("images.png", width=200)  # Replace with your logo file path

    # Apply custom CSS for styling
    st.markdown("""
    <style>
        .title {
            color: #4CAF50;
            font-size: 40px;
            font-weight: bold;
        }
        .button {
            background-color: #4CAF50;
            color: white;
        }
        .info {
            color: #2196F3;
        }
        .success {
            color: #4CAF50;
        }
    </style>
    """, unsafe_allow_html=True)

    st.markdown('<p class="title">PowerPoint Presentation Generator with Gemini API</p>', unsafe_allow_html=True)

    # Topic input and button
    topic = st.text_input("Enter the topic for your presentation:", key="topic_input")
    generate_button = st.button("Generate Presentation", key="generate_button", help="Click to generate your presentation", use_container_width=True)

    if generate_button and topic:
        st.markdown('<p class="info">Generating presentation... Please wait.</p>', unsafe_allow_html=True)
        
        # Generate slide titles and contents
        slide_titles = generate_slide_titles(topic)
        filtered_slide_titles = [item for item in slide_titles if item.strip() != '']
        slide_contents = [generate_slide_content(title) for title in filtered_slide_titles]
        create_presentation(topic, filtered_slide_titles, slide_contents)

        st.markdown('<p class="success">Presentation generated successfully!</p>', unsafe_allow_html=True)
        st.markdown(get_ppt_download_link(topic), unsafe_allow_html=True)

def get_ppt_download_link(topic):
    ppt_filename = f"generated_ppt/{topic}_presentation.pptx"

    with open(ppt_filename, "rb") as file:
        ppt_contents = file.read()

    b64_ppt = base64.b64encode(ppt_contents).decode()
    return f'<a href="data:application/vnd.openxmlformats-officedocument.presentationml.presentation;base64,{b64_ppt}" download="{topic}_presentation.pptx" style="color: #4CAF50; font-size: 18px;">Download the PowerPoint Presentation</a>'

if __name__ == "__main__":
    main()
